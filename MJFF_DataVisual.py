# import streamlit as st
# import pandas as pd
# import plotly.graph_objects as go
# from plotly.subplots import make_subplots
# import plotly.express as px
# import re
# from groq import Groq
# from io import BytesIO

# # --- PPTX GENERATION DEPENDENCIES ---
# try:
#     from pptx import Presentation
#     from pptx.util import Inches, Pt
#     from pptx.dml.color import RGBColor
#     from pptx.enum.text import PP_ALIGN
# except ImportError:
#     st.error("Missing dependency: please run `pip install python-pptx kaleido` in your terminal environment.")

# # --- INITIALIZATION ---
# if 'master_data' not in st.session_state:
#     st.session_state.master_data = pd.DataFrame(columns=[
#         'Include', 'Sample ID', 'Datestamp', 'Sample Name', 'Run', 'Description'
#     ])

# if 'raw_metrics_df' not in st.session_state:
#     st.session_state.raw_metrics_df = pd.DataFrame()

# if 'processed_files' not in st.session_state:
#     st.session_state.processed_files = set()

# if 'show_plotting' not in st.session_state:
#     st.session_state.show_plotting = False

# if 'num_plots' not in st.session_state:
#     st.session_state.num_plots = 1

# st.set_page_config(page_title="MJFF Sample Manager", layout="wide")

# # --- CUSTOM CSS FOR "WING" HINT ---
# st.markdown("""
#     <style>
#     @keyframes nudge {
#       0% { transform: translateX(0); }
#       50% { transform: translateX(5px); }
#       100% { transform: translateX(0); }
#     }
#     .insight-hint {
#         background-color: #e8f4f8;
#         border-left: 5px solid #007bff;
#         padding: 12px;
#         border-radius: 8px;
#         font-size: 0.95rem;
#         margin-bottom: 15px;
#         animation: nudge 3s ease-in-out 3;
#         display: flex;
#         align-items: center;
#         gap: 10px;
#     }
#     </style>
# """, unsafe_allow_html=True)

# st.title("MJFF Analysis Data Visualization")

# # --- APP FUNCTIONALITY GUIDE ---
# with st.expander("📖 How to use this App (Feature Guide)", expanded=True):
#     st.markdown("""
#     ### 🚀 Getting Started
#     1. **Upload Data**: Use the CSV uploader below. The app parses Sample IDs between boundaries (`_xcxcx-`, `-xzxzx-`, `_cxcx_`) and isolates pure text.
#     2. **Auto Run Counts**: Run sequences ($1, 2, 3, \dots$) are automatically computed based on sample name occurrence order during import.
    
#     ### 🛠️ Key Functionalities
#     * **Bulk Management**: Use **Section 1** tools to batch-rename samples, append structural elements, or re-index counts on demand.
#     * **Target Dropdown Filter**: In **Section 2**, cleanly isolate individual groups from your dropdown menu to instantly re-render active viewports.
#     * **📊 Cross-Metric Outlier Filtering**: Filter out anomalies using one parameter (e.g., *Barcode std.*) and view the surviving records plotted in your core parameter (e.g., *PP-Gauss*).
#     * **📈 Cohort Summaries & CV% Shifts**: Section 3 tracks quantitative and visual shifts in the Coefficient of Variation ($CV\%$) between raw and refined streams.
#     * **📝 Automated PPTX Generation**: Section 6 auto-compiles all charts and summary metrics into a structured presentation deck.
#     """)

# # --- AI INSIGHT ENGINE ---
# def get_detailed_insights(df, metric, lens_type, custom_query=None):
#     try:
#         client = Groq(api_key=st.secrets["GROQ_API_KEY"])
#     except KeyError:
#         return "Error: GROQ_API_KEY not found in Streamlit Secrets."
    
#     summary_report = []
#     other_params = ['Barcode av.', 'Barcode std.', 'NrDrops']
#     other_params = [p for p in other_params if p in df.columns]
    
#     for name in df['Sample Name'].unique():
#         sub = df[df['Sample Name'] == name]
#         mean = sub[metric].mean()
#         std = sub[metric].std()
#         cv = (std / mean * 100) if mean != 0 else 0
        
#         outliers = sub[(sub[metric] > mean + 1*std) | (sub[metric] < mean - 1*std)]
#         group_head = f"GROUP: {name}\n- Primary Metric ({metric}) Mean: {mean:.2f} (±{std:.2f})\n- Group CV: {cv:.1f}%"
        
#         outlier_notes = []
#         if not outliers.empty:
#             for _, row in outliers.iterrows():
#                 note = f"  * OUTLIER DETECTED: Run {row['Run']} value is {row[metric]:.2f}."
#                 evidence = []
#                 for p in other_params:
#                     p_mean = sub[p].mean()
#                     p_std = sub[p].std()
#                     if p_std > 0 and abs(row[p] - p_mean) > 1 * p_std:
#                         evidence.append(f"{p} is abnormal at {row[p]:.2f} (Group Avg: {p_mean:.2f})")
#                 note += "\n    EVIDENCE: " + (" | ".join(evidence) if evidence else "Secondary parameters were stable.")
#                 outlier_notes.append(note)
#         else:
#             outlier_notes.append("  * No outliers detected.")
#         summary_report.append(group_head + "\n" + "\n".join(outlier_notes))

#     full_data_summary = "\n\n".join(summary_report)
#     prompts = {
#         "General Analyst": "Summarize performance and outliers.",
#         "Quality Control Specialist": "Critique stability; suggest if runs should be discarded.",
#         "Root Cause Investigator": "Forensic focus. Correlate outliers with Barcode/Drops data."
#     }

#     system_msg = f"You are an expert {lens_type}. " + prompts[lens_type]
#     user_msg = f"LAB DATA ANALYSIS REQUEST:\nPRIMARY METRIC: {metric}\n\nDATA SUMMARY:\n{full_data_summary}"
    
#     if custom_query:
#         user_msg += f"\n\nUSER QUESTION: {custom_query}\n\nPlease prioritize answering the user's question using the data provided."

#     try:
#         completion = client.chat.completions.create(
#             model="llama-3.3-70b-versatile",
#             messages=[{"role": "system", "content": system_msg}, {"role": "user", "content": user_msg}],
#             temperature=0.1
#         )
#         return completion.choices[0].message.content
#     except Exception as e:
#         return f"AI Error: {str(e)}"

# # --- FILE UPLOADER WITH AUTOMATED SEQUENCE GENERATION ---
# uploaded_file = st.file_uploader("Upload CSV", type="csv")

# if uploaded_file:
#     if uploaded_file.name not in st.session_state.processed_files:
#         df = pd.read_csv(uploaded_file)
#         if 'Sample' in df.columns:
#             st.session_state.raw_metrics_df = pd.concat(
#                 [st.session_state.raw_metrics_df, df], ignore_index=True
#             ).drop_duplicates(subset=['Sample'])
            
#             parsing_list = []
#             for s_id in df['Sample'].unique():
#                 if s_id not in st.session_state.master_data['Sample ID'].tolist():
#                     s_id_str = str(s_id)
                    
#                     date_match = re.search(r'(\d{8})[-_]', s_id_str)
#                     ds = date_match.group(1) if date_match else "Unknown"
                    
#                     boundary_match = re.search(r'(?:[-_][a-zA-Z]+[-_])([a-zA-Z]+)', s_id_str)
                    
#                     if boundary_match:
#                         sample_name = boundary_match.group(1)
#                     else:
#                         rest = s_id_str.split(ds, 1)[1].lstrip('-_') if date_match else s_id_str
#                         sample_name = re.sub(r'[^a-zA-Z]', '', rest.split('_')[0])
                    
#                     parsing_list.append({
#                         'Include': True, 
#                         'Sample ID': s_id, 
#                         'Datestamp': ds, 
#                         'Sample Name': sample_name, 
#                         'Description': ""
#                     })
            
#             if parsing_list:
#                 new_batch_df = pd.DataFrame(parsing_list)
#                 combined_temp = pd.concat([st.session_state.master_data, new_batch_df], ignore_index=True)
#                 combined_temp['Run'] = (combined_temp.groupby('Sample Name').cumcount() + 1).astype(str)
#                 st.session_state.master_data = combined_temp
                
#             st.session_state.processed_files.add(uploaded_file.name)
#             st.rerun()

# if not st.session_state.master_data.empty:
#     st.subheader("1. Sample ID List Table")

#     # --- BULK TOOLS SECTION ---
#     with st.expander("🛠️ Bulk Rename & Selection Tools"):
#         b_col1, b_col2, b_col3, b_col4 = st.columns(4)
#         with b_col1:
#             st.markdown("**Find & Replace**")
#             f_txt = st.text_input("Text to find", key="bulk_find")
#             r_txt = st.text_input("Replace with", key="bulk_replace")
#             if st.button("Replace in Names", key="btn_replace"):
#                 st.session_state.master_data['Sample Name'] = st.session_state.master_data['Sample Name'].str.replace(f_txt, r_txt)
#                 st.session_state.master_data['Run'] = (st.session_state.master_data.groupby('Sample Name').cumcount() + 1).astype(str)
#                 st.rerun()
#         with b_col2:
#             st.markdown("**Add Text**")
#             pre_txt = st.text_input("Add Prefix", key="bulk_prefix")
#             suf_txt = st.text_input("Add Suffix", key="bulk_suffix")
#             if st.button("Apply Text", key="btn_apply_text"):
#                 if pre_txt: st.session_state.master_data['Sample Name'] = pre_txt + st.session_state.master_data['Sample Name']
#                 if suf_txt: st.session_state.master_data['Sample Name'] = st.session_state.master_data['Sample Name'] + suf_txt
#                 st.session_state.master_data['Run'] = (st.session_state.master_data.groupby('Sample Name').cumcount() + 1).astype(str)
#                 st.rerun()
#         with b_col3:
#             st.markdown("**Batch Run Update**")
#             mode = st.radio("Run Logic", ["Set Constant", "Recalculate Occurrences (1,2,3...)"], key="bulk_run_mode")
#             val = st.number_input("Start Value Offset", min_value=0, step=1, value=1, key="bulk_run_val")
#             if st.button("Update Runs", key="btn_update_runs"):
#                 if mode == "Set Constant":
#                     st.session_state.master_data['Run'] = str(val)
#                 else:
#                     st.session_state.master_data['Run'] = (
#                         st.session_state.master_data.groupby('Sample Name').cumcount() + val
#                     ).astype(str)
#                 st.rerun()
#         with b_col4:
#             st.markdown("**Selection Control**")
#             if st.button("✅ Select All", key="btn_select_all"):
#                 st.session_state.master_data['Include'] = True
#                 st.rerun()
#             if st.button("❌ Deselect All", key="btn_deselect_all"):
#                 st.session_state.master_data['Include'] = False
#                 st.rerun()

#     edited_df = st.data_editor(
#         st.session_state.master_data,
#         key="editor_widget",
#         column_config={
#             "Include": st.column_config.CheckboxColumn("Plot?", default=True),
#             "Sample ID": st.column_config.TextColumn("Original ID", disabled=True),
#             "Datestamp": st.column_config.TextColumn("Date", disabled=True),
#             "Run": st.column_config.TextColumn("Run Custom Override"),
#         },
#         hide_index=True,
#         width="stretch",
#     )

#     col1, col2 = st.columns([1, 5])
#     with col1:
#         if st.button("Commit Changes", key="btn_commit"):
#             st.session_state.master_data = edited_df
#             st.success("Changes saved!")
    
#     with col2:
#         if st.button("Proceed to Plotting 📊", key="btn_proceed"):
#             st.session_state.master_data = edited_df 
#             st.session_state.show_plotting = True

# # --- PLOTTING & INSIGHTS ---
# if st.session_state.show_plotting:
#     st.divider()
#     st.markdown('<div class="insight-hint"><span>🤖</span><b>Automated Insights Available:</b> Scroll to Section 5 for AI root-cause analysis and custom queries.</div>', unsafe_allow_html=True)
    
#     st.subheader("2. Comparison Plotting")
#     base_filtered = st.session_state.master_data[st.session_state.master_data['Include'] == True]

#     if base_filtered.empty:
#         st.warning("No samples marked for selection. Please check the 'Plot?' boxes in the table above.")
#     else:
#         # --- FIXED STICKY PALETTE INITIALIZATION ---
#         all_unique_names = sorted(st.session_state.master_data['Sample Name'].unique())
#         color_palette = px.colors.qualitative.Plotly
#         sticky_palette = {name: color_palette[idx % len(color_palette)] for idx, name in enumerate(all_unique_names)}

#         available_sample_names = sorted(base_filtered['Sample Name'].unique())
#         selected_samples = st.multiselect(
#             "🔎 Select Sample Names to Plot from Dropdown:",
#             options=available_sample_names,
#             default=available_sample_names,
#             key="sample_dropdown_filter"
#         )
        
#         final_filtered_master = base_filtered[base_filtered['Sample Name'].isin(selected_samples)]

#         if final_filtered_master.empty:
#             st.info("Please select at least one sample name from the dropdown menu to render plots.")
#         else:
#             plot_df = pd.merge(
#                 final_filtered_master, 
#                 st.session_state.raw_metrics_df, 
#                 left_on='Sample ID', 
#                 right_on='Sample'
#             )
            
#             core_metrics = ['PP-Gauss', 'PP-750', 'Barcode av.', 'Barcode std.', 'NrDrops']
#             available_metrics = [m for m in core_metrics if m in plot_df.columns]

#             # --- GLOBAL OUTLIER FILTER CONTROLS ---
#             st.markdown("#### 🛡️ Cross-Metric Outlier Configuration")
#             out_c1, out_c2, out_c3 = st.columns([1, 1, 2])
#             with out_c1:
#                 filter_metric = st.selectbox(
#                     "Metric to evaluate for Outliers:",
#                     options=available_metrics,
#                     index=available_metrics.index('Barcode std.') if 'Barcode std.' in available_metrics else 0,
#                     key="global_filter_metric"
#                 )
#             with out_c2:
#                 sigma_multiplier = st.selectbox(
#                     "Outlier Boundary (Sigma):",
#                     options=[1, 2, 3],
#                     index=1,
#                     format_func=lambda x: f"{x} Sigma ({x}σ)",
#                     key="global_sigma_multiplier"
#                 )
#             with out_c3:
#                 st.caption(f"**Rule Logic:** The system calculates the mean and standard deviation of **{filter_metric}** for each cohort. Points outside ±{sigma_multiplier}σ will be completely stripped out from the cleaned views.")

#             view_mode = st.radio(
#                 "Select Visualization Mode:",
#                 ["Detailed (Show Every Run)", "Global (Mean & Std Dev)"],
#                 horizontal=True,
#                 key="view_mode_selector"
#             )

#             # Pre-calculate the cleaned dataset globally so it matches across all components
#             cleaned_plot_df_list = []
#             for name in plot_df['Sample Name'].unique():
#                 sub = plot_df[plot_df['Sample Name'] == name]
#                 if len(sub) > 1:
#                     mean_val = sub[filter_metric].mean()
#                     std_val = sub[filter_metric].std()
#                     if pd.isna(std_val) or std_val == 0:
#                         cleaned_plot_df_list.append(sub)
#                     else:
#                         cutoff = std_val * sigma_multiplier
#                         filtered_sub = sub[
#                             (sub[filter_metric] >= (mean_val - cutoff)) & 
#                             (sub[filter_metric] <= (mean_val + cutoff))
#                         ]
#                         cleaned_plot_df_list.append(filtered_sub)
#                 else:
#                     cleaned_plot_df_list.append(sub)
            
#             global_cleaned_df = pd.concat(cleaned_plot_df_list, ignore_index=True) if cleaned_plot_df_list else plot_df.copy()

#             lines_removed = len(plot_df) - len(global_cleaned_df)
#             if lines_removed > 0:
#                 st.toast(f"✂️ Pruned {lines_removed} outlier runs based on {filter_metric} ({sigma_multiplier}σ)!", icon="ℹ️")

#             # --- LOOP THROUGH PLOT WINDOWS ---
#             for i in range(st.session_state.num_plots):
#                 st.markdown(f"---")
#                 st.markdown(f"### Plot Window {i+1}")
#                 c1, c2 = st.columns(2)
                
#                 with c1:
#                     selected_metric = st.selectbox(
#                         f"Select Plotting Metric", 
#                         available_metrics, 
#                         key=f"metric_select_{i}",
#                         index=available_metrics.index('PP-Gauss') if 'PP-Gauss' in available_metrics and i == 0 else i % len(available_metrics)
#                     )
#                 with c2:
#                     sort_order = st.selectbox(
#                         "Sort Order",
#                         ["Default (Name)", "Ascending (by Mean)", "Descending (by Mean)"],
#                         key=f"sort_select_{i}"
#                     )

#                 def process_and_render_data(data_to_plot, is_clean_view):
#                     if view_mode == "Detailed (Show Every Run)":
#                         temp_df = data_to_plot.copy()
                        
#                         if sort_order != "Default (Name)":
#                             group_means = temp_df.groupby('Sample Name')[selected_metric].mean()
#                             ascending = True if "Ascending" in sort_order else False
#                             sorted_names = group_means.sort_values(ascending=ascending).index
#                             temp_df['Sample Name'] = pd.Categorical(temp_df['Sample Name'], categories=sorted_names, ordered=True)
#                             temp_df = temp_df.sort_values(['Sample Name', 'Datestamp', 'Run'])
#                         else:
#                             temp_df = temp_df.sort_values(by=['Sample Name', 'Datestamp', 'Run'])

#                         temp_df['Plot_X'] = (
#                             "<span style='color:teal; font-weight:bold'>" + temp_df['Sample Name'].astype(str) + "</span><br>" + 
#                             "<span style='color:gray'>" + temp_df['Datestamp'].astype(str) + "</span><br>" + 
#                             "<span style='color:tomato'>Run: " + temp_df['Run'].astype(str) + "</span>"
#                         )
                        
#                         y_val = temp_df[selected_metric]
#                         error_val = None
#                         custom_data = temp_df[['Description']]
#                         htemp = "<b>%{x}</b><br>Value: %{y}<br>Description: %{customdata[0]}<extra></extra>"
#                         display_df = temp_df

#                     else:
#                         agg_results = data_to_plot.groupby('Sample Name')[available_metrics].agg(['mean', 'std']).reset_index()
#                         agg_results.columns = [f"{col[0]}_{col[1]}" if col[1] else col[0] for col in agg_results.columns]
                        
#                         if sort_order != "Default (Name)":
#                             ascending = True if "Ascending" in sort_order else False
#                             agg_results = agg_results.sort_values(by=f"{selected_metric}_mean", ascending=ascending)
                        
#                         agg_results['Plot_X'] = "<span style='color:teal; font-weight:bold'>" + agg_results['Sample Name'] + "</span>"
                        
#                         y_val = agg_results[f"{selected_metric}_mean"]
#                         error_val = agg_results[f"{selected_metric}_std"]
#                         custom_data = error_val
#                         htemp = "<b>%{x}</b><br>Mean: %{y:.2f}<br>Std Dev: %{customdata:.2f}<extra></extra>"
#                         display_df = agg_results

#                     display_df['BarColor'] = display_df['Sample Name'].map(sticky_palette)

#                     fig = go.Figure()
#                     fig.add_trace(
#                         go.Bar(
#                             x=display_df['Plot_X'], 
#                             y=y_val, 
#                             marker_color=display_df['BarColor'],
#                             text=y_val.round(2) if not y_val.isna().all() else "",
#                             textposition='auto',
#                             error_y=dict(type='data', array=error_val, visible=True) if view_mode == "Global (Mean & Std Dev)" else None,
#                             customdata=custom_data,
#                             hovertemplate=htemp,
#                             textfont=dict(size=13, weight="bold")
#                         )
#                     )

#                     title_lbl = f"{selected_metric} (Cleaned via {filter_metric} @ {sigma_multiplier}σ)" if is_clean_view else f"{selected_metric} (Unfiltered Original Data)"
                    
#                     fig.update_layout(
#                         title=dict(text=title_lbl, font=dict(size=18, weight="bold")),
#                         height=480, 
#                         yaxis_title=selected_metric, 
#                         template="plotly_white", 
#                         showlegend=False,
#                         yaxis=dict(
#                             range=[0, y_val.max() * 1.2] if not y_val.empty and not y_val.isna().all() else [0, 1],
#                             title_font=dict(size=14, weight="bold"),
#                             tickfont=dict(size=12)
#                         ),
#                         xaxis=dict(tickfont=dict(size=12)),
#                         hoverlabel=dict(font_size=14)
#                     )
#                     fig.update_xaxes(tickangle=0)
#                     return fig

#                 graph_col1, graph_col2 = st.columns(2)
#                 with graph_col1:
#                     st.plotly_chart(process_and_render_data(plot_df, is_clean_view=False), width="stretch", key=f"chart_orig_{i}")
#                 with graph_col2:
#                     st.plotly_chart(process_and_render_data(global_cleaned_df, is_clean_view=True), width="stretch", key=f"chart_clean_{i}")

#             col_btn1, col_btn2 = st.columns([1, 4])
#             with col_btn1:
#                 if st.button("➕ Add Plot", key="btn_add_plot"):
#                     st.session_state.num_plots += 1
#                     st.rerun()
#             with col_btn2:
#                 if st.session_state.num_plots > 1:
#                     if st.button("Reset Plots", key="btn_reset_plots"):
#                         st.session_state.num_plots = 1
#                         st.rerun()

#             # --- SECTION 3: COHORT SUMMARY STATISTICS & CV% COMPARISON ---
#             st.write("---")
#             st.subheader("3. Cohort Summary Statistics Table & CV% Stability Delta")
            
#             has_gauss = 'PP-Gauss' in plot_df.columns
#             has_750 = 'PP-750' in plot_df.columns

#             if (has_gauss or has_750) and not plot_df.empty:
                
#                 def generate_summary_metrics(dataframe, suffix):
#                     agg_dict = {}
#                     if has_gauss: agg_dict['PP-Gauss'] = ['count', 'mean', 'std']
#                     if has_750: agg_dict['PP-750'] = ['count', 'mean', 'std']
                    
#                     summary = dataframe.groupby('Sample Name').agg(agg_dict).reset_index()
#                     summary.columns = [f"{col[0]}_{col[1]}" if col[1] else col[0] for col in summary.columns]
                    
#                     final_cols = ['Sample Name']
#                     if has_gauss:
#                         summary[f'Total_Runs_Gauss_{suffix}'] = summary['PP-Gauss_count']
#                         summary[f'Mean_Gauss_{suffix}'] = summary['PP-Gauss_mean']
#                         summary[f'Std_Gauss_{suffix}'] = summary['PP-Gauss_std']
#                         summary[f'CV%_Gauss_{suffix}'] = (summary[f'Std_Gauss_{suffix}'] / summary[f'Mean_Gauss_{suffix}']) * 100
#                         final_cols.extend([f'Total_Runs_Gauss_{suffix}', f'Mean_Gauss_{suffix}', f'Std_Gauss_{suffix}', f'CV%_Gauss_{suffix}'])
                    
#                     if has_750:
#                         summary[f'Total_Runs_750_{suffix}'] = summary['PP-750_count']
#                         summary[f'Mean_750_{suffix}'] = summary['PP-750_mean']
#                         summary[f'Std_750_{suffix}'] = summary['PP-750_std']
#                         summary[f'CV%_750_{suffix}'] = (summary[f'Std_750_{suffix}'] / summary[f'Mean_750_{suffix}']) * 100
#                         final_cols.extend([f'Total_Runs_750_{suffix}', f'Mean_750_{suffix}', f'Std_750_{suffix}', f'CV%_750_{suffix}'])
                        
#                     return summary[final_cols].fillna(0)

#                 orig_metrics = generate_summary_metrics(plot_df, "orig")
#                 clean_metrics = generate_summary_metrics(global_cleaned_df, "clean")
#                 merged_summary = pd.merge(orig_metrics, clean_metrics, on='Sample Name', how='outer').fillna(0)

#                 # --- CV% COMPARISON PLOT ---
#                 if has_gauss:
#                     cv_fig = go.Figure()
#                     cv_fig.add_trace(go.Bar(
#                         x=merged_summary['Sample Name'],
#                         y=merged_summary['CV%_Gauss_orig'],
#                         name='Original Dataset PP-Gauss CV%',
#                         marker_color='#ef553b',
#                         text=merged_summary['CV%_Gauss_orig'].round(2).astype(str) + '%',
#                         textposition='auto',
#                         textfont=dict(size=12, weight="bold")
#                     ))
#                     cv_fig.add_trace(go.Bar(
#                         x=merged_summary['Sample Name'],
#                         y=merged_summary['CV%_Gauss_clean'],
#                         name=f'Cleaned Dataset PP-Gauss CV% (via {filter_metric})',
#                         marker_color='#636efa',
#                         text=merged_summary['CV%_Gauss_clean'].round(2).astype(str) + '%',
#                         textposition='auto',
#                         textfont=dict(size=12, weight="bold")
#                     ))
#                     cv_fig.update_layout(
#                         title=dict(text=f"Stability Impact: PP-Gauss CV% Delta After Outlier Pruning", font=dict(size=18, weight="bold")),
#                         xaxis_title="Sample Cohorts",
#                         yaxis_title="Coefficient of Variation (CV %)",
#                         barmode='group',
#                         template='plotly_white',
#                         height=420,
#                         xaxis=dict(title_font=dict(size=14, weight="bold"), tickfont=dict(size=12)),
#                         yaxis=dict(title_font=dict(size=14, weight="bold"), tickfont=dict(size=12)),
#                         legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1, font=dict(size=12))
#                     )
#                     st.plotly_chart(cv_fig, width="stretch", key="cv_comparison_chart")

#                 # --- INTERACTIVE BREAKDOWN SUB-TABLE ---
#                 sum_c1, sum_c2 = st.columns([1, 3])
#                 with sum_c1:
#                     summary_data_source = st.radio(
#                         "Summary Metric Source for Table:",
#                         ["Use Cleaned Dataset", "Use Original Dataset"],
#                         key="summary_dataset_toggle"
#                     )
#                 with sum_c2:
#                     st.caption(f"Granular analytical cross-table showing tracking indices and stability benchmarks for **PP-Gauss** and **PP-750** simultaneously.")

#                 if summary_data_source == "Use Cleaned Dataset":
#                     table_summary_view = clean_metrics.rename(columns={
#                         'Total_Runs_Gauss_clean': 'Total_Runs',
#                         'Mean_Gauss_clean': 'Mean_PP_Gauss', 'Std_Gauss_clean': 'Std_PP_Gauss', 'CV%_Gauss_clean': 'CV%_Gauss',
#                         'Mean_750_clean': 'Mean_PP_750', 'Std_750_clean': 'Std_PP_750', 'CV%_750_clean': 'CV%_750'
#                     })
#                 else:
#                     table_summary_view = orig_metrics.rename(columns={
#                         'Total_Runs_Gauss_orig': 'Total_Runs',
#                         'Mean_Gauss_orig': 'Mean_PP_Gauss', 'Std_Gauss_orig': 'Std_PP_Gauss', 'CV%_Gauss_orig': 'CV%_Gauss',
#                         'Mean_750_orig': 'Mean_PP_750', 'Std_750_orig': 'Std_PP_750', 'CV%_750_orig': 'CV%_750'
#                     })

#                 st.dataframe(
#                     table_summary_view[['Sample Name', 'Total_Runs', 'Mean_PP_Gauss', 'Std_PP_Gauss', 'CV%_Gauss', 'Mean_PP_750', 'Std_PP_750', 'CV%_750']],
#                     column_config={
#                         "Sample Name": st.column_config.TextColumn("Sample Name"),
#                         "Total_Runs": st.column_config.NumberColumn("Total Runs", format="%d"),
#                         "Mean_PP_Gauss": st.column_config.NumberColumn("Mean PP-Gauss", format="%.3f"),
#                         "Std_PP_Gauss": st.column_config.NumberColumn("Std. Dev PP-Gauss", format="%.3f"),
#                         "CV%_Gauss": st.column_config.NumberColumn("CV% (Gauss)", format="%.2f%%"),
#                         "Mean_PP_750": st.column_config.NumberColumn("Mean PP-750", format="%.3f"),
#                         "Std_PP_750": st.column_config.NumberColumn("Std. Dev PP-750", format="%.3f"),
#                         "CV%_750": st.column_config.NumberColumn("CV% (750)", format="%.2f%%"),
#                     },
#                     hide_index=True,
#                     width="stretch"
#                 )
#             else:
#                 st.info("Metrics not found or dataset view empty. Summary table skipped.")

#             # --- SECTION 4: PLOTTED DATA REFERENCE TABLE ---
#             st.write("---")
#             st.subheader("4. Plotted Data Reference Table")
            
#             t_cfg1, t_cfg2 = st.columns([1, 3])
#             with t_cfg1:
#                 table_view_mode = st.radio(
#                     "Table Data Source:",
#                     ["Show Cleaned Dataset", "Show Original Unfiltered Dataset"],
#                     key="table_view_mode"
#                 )
#             with t_cfg2:
#                 if table_view_mode == "Show Cleaned Dataset":
#                     st.caption(f"Showing **{len(global_cleaned_df)} surviving rows**. Outliers based on variance in `{filter_metric}` have been filtered out.")
#                 else:
#                     st.caption(f"Showing **{len(plot_df)} complete original rows**.")

#             selected_table_source = global_cleaned_df if table_view_mode == "Show Cleaned Dataset" else plot_df
#             table_df = selected_table_source.copy()
            
#             if not table_df.empty:
#                 table_df['SAMPLE ID'] = table_df['Datestamp'].astype(str) + "_" + table_df['Sample Name'].astype(str) + "_" + table_df['Run'].astype(str)
                
#                 st.dataframe(
#                     table_df[['Datestamp', 'Sample Name', 'Run', 'SAMPLE ID'] + available_metrics],
#                     column_config={
#                         "Datestamp": st.column_config.TextColumn("Date", width=100),
#                         "Run": st.column_config.TextColumn("Run"), 
#                         "SAMPLE ID": st.column_config.TextColumn("Combined ID", width=300),
#                         **{m: st.column_config.NumberColumn(m, width=120, format="%.2f") for m in available_metrics}
#                     },
#                     hide_index=True,
#                     width="stretch"
#                 )
#             else:
#                 st.info("The selected dataset view is empty.")

#         # --- SECTION 5: AUTOMATED DATA INSIGHTS ---
#         st.divider()
#         st.subheader("5. 🤖 Automated Data Insights")
        
#         with st.expander("⚙️ AI Analysis & Custom Questions", expanded=True):
#             i_col1, i_col2, i_col3 = st.columns([1, 1, 1])
#             with i_col1:
#                 ai_metric = st.selectbox("Target Analysis Metric", available_metrics, key="ai_target_metric")
#                 lens = st.selectbox("Select Expert Lens", ["General Analyst", "Quality Control Specialist", "Root Cause Investigator"], key="ai_lens")
#             with i_col2:
#                 custom_q = st.text_area("💬 Human Language Query", placeholder="e.g., 'Is there a correlation between NrDrops and the outliers in Sample B?'", key="ai_query")
#             with i_col3:
#                 st.write(" ")
#                 st.write(" ")
#                 generate_btn = st.button("Generate Insights 🪄", width="stretch", key="btn_generate_insights")

#         if generate_btn:
#             with st.spinner("Analyzing..."):
#                 final_report = get_detailed_insights(plot_df, ai_metric, lens, custom_q)
#                 st.info(final_report)

#         # --- SECTION 6: AUTOMATED POWERPOINT REPORT ENGINE ---
#         st.divider()
#         st.subheader("6. 📝 PowerPoint Summary Report Engine")
#         st.caption("Generate a fully structured executive slide deck incorporating global metrics, cohort comparisons, and dedicated structural run slides per unique sample.")
        
#         if st.button("Build PowerPoint Presentation Summary 🚀", width="stretch", key="btn_build_pptx"):
#             with st.spinner("Rendering vector charts and drawing presentation canvas layers..."):
#                 try:
#                     # 1. Initialize Blank PowerPoint File Layout (Widescreen 16:9 Aspect Ratio)
#                     prs = Presentation()
#                     prs.slide_width = Inches(13.333)
#                     prs.slide_height = Inches(7.5)
#                     blank_layout = prs.slide_layouts[6] # Blank Slide Template

#                     # --- HELPER FUNCTION: ADD BASIC HEADER TO SLIDES ---
#                     def add_slide_header(slide, title_text):
#                         txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.8))
#                         tf = txBox.text_frame
#                         tf.word_wrap = True
#                         p = tf.paragraphs[0]
#                         p.text = title_text
#                         p.font.size = Pt(28)
#                         p.font.bold = True
#                         p.font.color.rgb = RGBColor(0, 51, 102)

#                     # --- SLIDE 1: DECK COVER TITLE SLIDE ---
#                     slide1 = prs.slides.add_slide(blank_layout)
#                     title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(2.5))
#                     tf = title_box.text_frame
#                     p1 = tf.paragraphs[0]
#                     p1.text = "MJFF Analytics Executive Report"
#                     p1.font.size = Pt(44)
#                     p1.font.bold = True
#                     p1.font.color.rgb = RGBColor(0, 51, 102)
                    
#                     p2 = tf.add_paragraph()
#                     p2.text = f"Automated Instrumentation Run Summary • Outlier Baseline: {filter_metric} ({sigma_multiplier}σ)"
#                     p2.font.size = Pt(18)
#                     p2.font.color.rgb = RGBColor(102, 102, 102)

#                     # --- SLIDE 2: COHORT STABILITY DELTA GRAPH ---
#                     if has_gauss:
#                         slide2 = prs.slides.add_slide(blank_layout)
#                         add_slide_header(slide2, "Stability Performance: PP-Gauss CV% Variations")
                        
#                         img_stream = BytesIO()
#                         cv_fig.write_image(img_stream, format="png", width=1100, height=500, scale=2)
#                         img_stream.seek(0)
#                         slide2.shapes.add_picture(img_stream, Inches(0.666), Inches(1.3), width=Inches(12.0))

#                     # --- SLIDE 3: COHORT DATA TABLE SUMMARY ---
#                     slide3 = prs.slides.add_slide(blank_layout)
#                     add_slide_header(slide3, f"Cohort Parameter Analytics View ({summary_data_source})")
                    
#                     # Force dataframe alignment to safe structural columns
#                     ppt_table_df = table_summary_view[['Sample Name', 'Total_Runs', 'Mean_PP_Gauss', 'Std_PP_Gauss', 'CV%_Gauss', 'Mean_PP_750', 'Std_PP_750', 'CV%_750']].reset_index(drop=True)
                    
#                     rows, cols = len(ppt_table_df) + 1, len(ppt_table_df.columns)
#                     left, top, width, height = Inches(0.5), Inches(1.5), Inches(12.333), Inches(0.4 * rows)
#                     table_shape = slide3.shapes.add_table(rows, cols, left, top, width, height)
#                     table = table_shape.table

#                     # Safe Display Headers Array matching exactly the column order above
#                     headers_display = ["Sample Name", "Total Runs", "Mean Gauss", "Std Dev Gauss", "CV% Gauss", "Mean 750", "Std Dev 750", "CV% 750"]
#                     for col_idx, text in enumerate(headers_display):
#                         cell = table.cell(0, col_idx)
#                         cell.text = text
#                         cell.fill.solid()
#                         cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
#                         p = cell.text_frame.paragraphs[0]
#                         p.font.bold = True
#                         p.font.color.rgb = RGBColor(255, 255, 255)
#                         p.font.size = Pt(12)
#                         p.alignment = PP_ALIGN.CENTER

#                     # Safe value population via strict position lookup
#                     for sequential_row_idx in range(len(ppt_table_df)):
#                         row_series = ppt_table_df.iloc[sequential_row_idx]
#                         for col_idx in range(len(row_series)):
#                             cell = table.cell(sequential_row_idx + 1, col_idx)
#                             value = row_series.iloc[col_idx]
                            
#                             # Clean cell formatting based on target metric name
#                             col_name = ppt_table_df.columns[col_idx]
#                             if isinstance(value, float):
#                                 cell.text = f"{value:.3f}%" if "CV%" in col_name else f"{value:.3f}"
#                             else:
#                                 cell.text = str(value)
                                
#                             p = cell.text_frame.paragraphs[0]
#                             p.font.size = Pt(11)
#                             p.alignment = PP_ALIGN.CENTER

#                     # --- SLIDES 4+: INDIVIDUAL SAMPLE PROFILE STACKS ---
#                     def build_micro_timeline(dataframe, sample_name, y_metric, title_text):
#                         sub_sub = dataframe[dataframe['Sample Name'] == sample_name].sort_values(by=['Datestamp', 'Run'])
#                         sub_sub['X_Label'] = "Run: " + sub_sub['Run'].astype(str)
                        
#                         fig_micro = go.Figure(data=[
#                             go.Bar(
#                                 x=sub_sub['X_Label'],
#                                 y=sub_sub[y_metric],
#                                 marker_color=sticky_palette.get(sample_name, '#636efa'),
#                                 text=sub_sub[y_metric].round(2) if not sub_sub[y_metric].isna().all() else "",
#                                 textposition='auto'
#                             )
#                         ])
#                         fig_micro.update_layout(
#                             title=dict(text=title_text, font=dict(size=14, weight="bold")),
#                             margin=dict(l=40, r=20, t=40, b=30),
#                             template="plotly_white",
#                             height=250,
#                             yaxis=dict(title=y_metric, title_font=dict(size=10), tickfont=dict(size=9)),
#                             xaxis=dict(tickfont=dict(size=10))
#                         )
#                         img_buf = BytesIO()
#                         fig_micro.write_image(img_buf, format="png", width=1100, height=220, scale=2)
#                         img_buf.seek(0)
#                         return img_buf

#                     unique_cohorts = sorted(plot_df['Sample Name'].unique())
                    
#                     for cohort in unique_cohorts:
#                         slide_c = prs.slides.add_slide(blank_layout)
#                         add_slide_header(slide_c, f"Sample Run Analysis: {cohort}")
                        
#                         if 'PP-Gauss' in plot_df.columns:
#                             img_g = build_micro_timeline(global_cleaned_df, cohort, 'PP-Gauss', "PP-Gauss Cleaned Signal Timeline")
#                             slide_c.shapes.add_picture(img_g, Inches(0.666), Inches(1.2), width=Inches(12.0))
                            
#                         if 'PP-750' in plot_df.columns:
#                             img_750 = build_micro_timeline(global_cleaned_df, cohort, 'PP-750', "PP-750 Cleaned Signal Timeline")
#                             slide_c.shapes.add_picture(img_750, Inches(0.666), Inches(3.2), width=Inches(12.0))
                            
#                         if 'Barcode av.' in plot_df.columns:
#                             img_bav = build_micro_timeline(global_cleaned_df, cohort, 'Barcode av.', "Barcode av. Reference Run Timeline")
#                             slide_c.shapes.add_picture(img_bav, Inches(0.666), Inches(5.2), width=Inches(12.0))

#                     # 4. Save and Deliver Presentation File Stream
#                     ppt_out = BytesIO()
#                     prs.save(ppt_out)
#                     ppt_out.seek(0)
                    
#                     st.success("PowerPoint compilation complete! Click below to download.")
#                     st.download_button(
#                         label="📥 Download PowerPoint Report (.pptx)",
#                         data=ppt_out,
#                         file_name="MJFF_Instrumentation_Summary_Report.pptx",
#                         mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
#                     )
#                 except Exception as e:
#                     st.error(f"Failed to generate presentation deck: {str(e)}")

import streamlit as st
import pandas as pd
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import plotly.express as px
import re
from groq import Groq
from io import BytesIO

# --- PPTX GENERATION DEPENDENCIES ---
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    st.error("Missing dependency: please run `pip install python-pptx kaleido` in your terminal environment.")

# --- INITIALIZATION ---
if 'master_data' not in st.session_state:
    st.session_state.master_data = pd.DataFrame(columns=[
        'Include', 'Sample ID', 'Datestamp', 'Sample Name', 'Run', 'Description'
    ])

if 'raw_metrics_df' not in st.session_state:
    st.session_state.raw_metrics_df = pd.DataFrame()

if 'processed_files' not in st.session_state:
    st.session_state.processed_files = set()

if 'show_plotting' not in st.session_state:
    st.session_state.show_plotting = False

if 'num_plots' not in st.session_state:
    st.session_state.num_plots = 1

st.set_page_config(page_title="MJFF Sample Manager", layout="wide")

# --- CUSTOM CSS FOR "WING" HINT ---
st.markdown("""
    <style>
    @keyframes nudge {
      0% { transform: translateX(0); }
      50% { transform: translateX(5px); }
      100% { transform: translateX(0); }
    }
    .insight-hint {
        background-color: #e8f4f8;
        border-left: 5px solid #007bff;
        padding: 12px;
        border-radius: 8px;
        font-size: 0.95rem;
        margin-bottom: 15px;
        animation: nudge 3s ease-in-out 3;
        display: flex;
        align-items: center;
        gap: 10px;
    }
    </style>
""", unsafe_allow_html=True)

st.title("MJFF Analysis Data Visualization")

# --- APP FUNCTIONALITY GUIDE ---
with st.expander("📖 How to use this App (Feature Guide)", expanded=True):
    st.markdown("""
    ### 🚀 Getting Started
    1. **Upload Data**: Use the CSV uploader below. The app parses Sample IDs between boundaries (`_xcxcx-`, `-xzxzx-`, `_cxcx_`) and isolates pure text.
    2. **Auto Run Counts**: Run sequences ($1, 2, 3, \dots$) are automatically computed based on sample name occurrence order during import.
    
    ### 🛠️ Key Functionalities
    * **Bulk Management**: Use **Section 1** tools to batch-rename samples, append structural elements, or re-index counts on demand.
    * **Target Dropdown Filter**: In **Section 2**, cleanly isolate individual groups from your dropdown menu to instantly re-render active viewports.
    * **📊 Cross-Metric Outlier Filtering**: Filter out anomalies using one parameter (e.g., *Barcode std.*) and view the surviving records plotted in your core parameter (e.g., *PP-Gauss*).
    * **📈 Cohort Summaries & CV% Shifts**: Section 3 tracks quantitative and visual shifts in the Coefficient of Variation ($CV\%$) between raw and refined streams.
    * **📝 Automated PPTX Generation**: Section 6 auto-compiles all macro system metrics, side-by-side comparative diagnostics, and individual data stacks.
    """)

# --- AI INSIGHT ENGINE ---
def get_detailed_insights(df, metric, lens_type, custom_query=None):
    try:
        client = Groq(api_key=st.secrets["GROQ_API_KEY"])
    except KeyError:
        return "Error: GROQ_API_KEY not found in Streamlit Secrets."
    
    summary_report = []
    other_params = ['Barcode av.', 'Barcode std.', 'NrDrops']
    other_params = [p for p in other_params if p in df.columns]
    
    for name in df['Sample Name'].unique():
        sub = df[df['Sample Name'] == name]
        mean = sub[metric].mean()
        std = sub[metric].std()
        cv = (std / mean * 100) if mean != 0 else 0
        
        outliers = sub[(sub[metric] > mean + 1*std) | (sub[metric] < mean - 1*std)]
        group_head = f"GROUP: {name}\n- Primary Metric ({metric}) Mean: {mean:.2f} (±{std:.2f})\n- Group CV: {cv:.1f}%"
        
        outlier_notes = []
        if not outliers.empty:
            for _, row in outliers.iterrows():
                note = f"  * OUTLIER DETECTED: Run {row['Run']} value is {row[metric]:.2f}."
                evidence = []
                for p in other_params:
                    p_mean = sub[p].mean()
                    p_std = sub[p].std()
                    if p_std > 0 and abs(row[p] - p_mean) > 1 * p_std:
                        evidence.append(f"{p} is abnormal at {row[p]:.2f} (Group Avg: {p_mean:.2f})")
                note += "\n    EVIDENCE: " + (" | ".join(evidence) if evidence else "Secondary parameters were stable.")
                outlier_notes.append(note)
        else:
            outlier_notes.append("  * No outliers detected.")
        summary_report.append(group_head + "\n" + "\n".join(outlier_notes))

    full_data_summary = "\n\n".join(summary_report)
    prompts = {
        "General Analyst": "Summarize performance and outliers.",
        "Quality Control Specialist": "Critique stability; suggest if runs should be discarded.",
        "Root Cause Investigator": "Forensic focus. Correlate outliers with Barcode/Drops data."
    }

    system_msg = f"You are an expert {lens_type}. " + prompts[lens_type]
    user_msg = f"LAB DATA ANALYSIS REQUEST:\nPRIMARY METRIC: {metric}\n\nDATA SUMMARY:\n{full_data_summary}"
    
    if custom_query:
        user_msg += f"\n\nUSER QUESTION: {custom_query}\n\nPlease prioritize answering the user's question using the data provided."

    try:
        completion = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[{"role": "system", "content": system_msg}, {"role": "user", "content": user_msg}],
            temperature=0.1
        )
        return completion.choices[0].message.content
    except Exception as e:
        return f"AI Error: {str(e)}"

# --- FILE UPLOADER WITH AUTOMATED SEQUENCE GENERATION ---
uploaded_file = st.file_uploader("Upload CSV", type="csv")

if uploaded_file:
    if uploaded_file.name not in st.session_state.processed_files:
        df = pd.read_csv(uploaded_file)
        if 'Sample' in df.columns:
            st.session_state.raw_metrics_df = pd.concat(
                [st.session_state.raw_metrics_df, df], ignore_index=True
            ).drop_duplicates(subset=['Sample'])
            
            parsing_list = []
            for s_id in df['Sample'].unique():
                if s_id not in st.session_state.master_data['Sample ID'].tolist():
                    s_id_str = str(s_id)
                    
                    date_match = re.search(r'(\d{8})[-_]', s_id_str)
                    ds = date_match.group(1) if date_match else "Unknown"
                    
                    boundary_match = re.search(r'(?:[-_][a-zA-Z]+[-_])([a-zA-Z]+)', s_id_str)
                    
                    if boundary_match:
                        sample_name = boundary_match.group(1)
                    else:
                        rest = s_id_str.split(ds, 1)[1].lstrip('-_') if date_match else s_id_str
                        sample_name = re.sub(r'[^a-zA-Z]', '', rest.split('_')[0])
                    
                    parsing_list.append({
                        'Include': True, 
                        'Sample ID': s_id, 
                        'Datestamp': ds, 
                        'Sample Name': sample_name, 
                        'Description': ""
                    })
            
            if parsing_list:
                new_batch_df = pd.DataFrame(parsing_list)
                combined_temp = pd.concat([st.session_state.master_data, new_batch_df], ignore_index=True)
                combined_temp['Run'] = (combined_temp.groupby('Sample Name').cumcount() + 1).astype(str)
                st.session_state.master_data = combined_temp
                
            st.session_state.processed_files.add(uploaded_file.name)
            st.rerun()

if not st.session_state.master_data.empty:
    st.subheader("1. Sample ID List Table")

    # --- BULK TOOLS SECTION ---
    with st.expander("🛠️ Bulk Rename & Selection Tools"):
        b_col1, b_col2, b_col3, b_col4 = st.columns(4)
        with b_col1:
            st.markdown("**Find & Replace**")
            f_txt = st.text_input("Text to find", key="bulk_find")
            r_txt = st.text_input("Replace with", key="bulk_replace")
            if st.button("Replace in Names", key="btn_replace"):
                st.session_state.master_data['Sample Name'] = st.session_state.master_data['Sample Name'].str.replace(f_txt, r_txt)
                st.session_state.master_data['Run'] = (st.session_state.master_data.groupby('Sample Name').cumcount() + 1).astype(str)
                st.rerun()
        with b_col2:
            st.markdown("**Add Text**")
            pre_txt = st.text_input("Add Prefix", key="bulk_prefix")
            suf_txt = st.text_input("Add Suffix", key="bulk_suffix")
            if st.button("Apply Text", key="btn_apply_text"):
                if pre_txt: st.session_state.master_data['Sample Name'] = pre_txt + st.session_state.master_data['Sample Name']
                if suf_txt: st.session_state.master_data['Sample Name'] = st.session_state.master_data['Sample Name'] + suf_txt
                st.session_state.master_data['Run'] = (st.session_state.master_data.groupby('Sample Name').cumcount() + 1).astype(str)
                st.rerun()
        with b_col3:
            st.markdown("**Batch Run Update**")
            mode = st.radio("Run Logic", ["Set Constant", "Recalculate Occurrences (1,2,3...)"], key="bulk_run_mode")
            val = st.number_input("Start Value Offset", min_value=0, step=1, value=1, key="bulk_run_val")
            if st.button("Update Runs", key="btn_update_runs"):
                if mode == "Set Constant":
                    st.session_state.master_data['Run'] = str(val)
                else:
                    st.session_state.master_data['Run'] = (
                        st.session_state.master_data.groupby('Sample Name').cumcount() + val
                    ).astype(str)
                st.rerun()
        with b_col4:
            st.markdown("**Selection Control**")
            if st.button("✅ Select All", key="btn_select_all"):
                st.session_state.master_data['Include'] = True
                st.rerun()
            if st.button("❌ Deselect All", key="btn_deselect_all"):
                st.session_state.master_data['Include'] = False
                st.rerun()

    edited_df = st.data_editor(
        st.session_state.master_data,
        key="editor_widget",
        column_config={
            "Include": st.column_config.CheckboxColumn("Plot?", default=True),
            "Sample ID": st.column_config.TextColumn("Original ID", disabled=True),
            "Datestamp": st.column_config.TextColumn("Date", disabled=True),
            "Run": st.column_config.TextColumn("Run Custom Override"),
        },
        hide_index=True,
        width="stretch",
    )

    col1, col2 = st.columns([1, 5])
    with col1:
        if st.button("Commit Changes", key="btn_commit"):
            st.session_state.master_data = edited_df
            st.success("Changes saved!")
    
    with col2:
        if st.button("Proceed to Plotting 📊", key="btn_proceed"):
            st.session_state.master_data = edited_df 
            st.session_state.show_plotting = True

# --- PLOTTING & INSIGHTS ---
if st.session_state.show_plotting:
    st.divider()
    st.markdown('<div class="insight-hint"><span>🤖</span><b>Automated Insights Available:</b> Scroll to Section 5 for AI root-cause analysis and custom queries.</div>', unsafe_allow_html=True)
    
    st.subheader("2. Comparison Plotting")
    base_filtered = st.session_state.master_data[st.session_state.master_data['Include'] == True]

    if base_filtered.empty:
        st.warning("No samples marked for selection. Please check the 'Plot?' boxes in the table above.")
    else:
        # --- FIXED STICKY PALETTE INITIALIZATION ---
        all_unique_names = sorted(st.session_state.master_data['Sample Name'].unique())
        color_palette = px.colors.qualitative.Plotly
        sticky_palette = {name: color_palette[idx % len(color_palette)] for idx, name in enumerate(all_unique_names)}

        available_sample_names = sorted(base_filtered['Sample Name'].unique())
        selected_samples = st.multiselect(
            "🔎 Select Sample Names to Plot from Dropdown:",
            options=available_sample_names,
            default=available_sample_names,
            key="sample_dropdown_filter"
        )
        
        final_filtered_master = base_filtered[base_filtered['Sample Name'].isin(selected_samples)]

        if final_filtered_master.empty:
            st.info("Please select at least one sample name from the dropdown menu to render plots.")
        else:
            plot_df = pd.merge(
                final_filtered_master, 
                st.session_state.raw_metrics_df, 
                left_on='Sample ID', 
                right_on='Sample'
            )
            
            core_metrics = ['PP-Gauss', 'PP-750', 'Barcode av.', 'Barcode std.', 'NrDrops']
            available_metrics = [m for m in core_metrics if m in plot_df.columns]

            # --- GLOBAL OUTLIER FILTER CONTROLS ---
            st.markdown("#### 🛡️ Cross-Metric Outlier Configuration")
            out_c1, out_c2, out_c3 = st.columns([1, 1, 2])
            with out_c1:
                filter_metric = st.selectbox(
                    "Metric to evaluate for Outliers:",
                    options=available_metrics,
                    index=available_metrics.index('Barcode std.') if 'Barcode std.' in available_metrics else 0,
                    key="global_filter_metric"
                )
            with out_c2:
                sigma_multiplier = st.selectbox(
                    "Outlier Boundary (Sigma):",
                    options=[1, 2, 3],
                    index=1,
                    format_func=lambda x: f"{x} Sigma ({x}σ)",
                    key="global_sigma_multiplier"
                )
            with out_c3:
                st.caption(f"**Rule Logic:** The system calculates the mean and standard deviation of **{filter_metric}** for each cohort. Points outside $\pm${sigma_multiplier}$\sigma$ will be completely stripped out from the cleaned views.")

            view_mode = st.radio(
                "Select Visualization Mode:",
                ["Detailed (Show Every Run)", "Global (Mean & Std Dev)"],
                horizontal=True,
                key="view_mode_selector"
            )

            # Pre-calculate the cleaned dataset globally so it matches across all components
            cleaned_plot_df_list = []
            for name in plot_df['Sample Name'].unique():
                sub = plot_df[plot_df['Sample Name'] == name]
                if len(sub) > 1:
                    mean_val = sub[filter_metric].mean()
                    std_val = sub[filter_metric].std()
                    if pd.isna(std_val) or std_val == 0:
                        cleaned_plot_df_list.append(sub)
                    else:
                        cutoff = std_val * sigma_multiplier
                        filtered_sub = sub[
                            (sub[filter_metric] >= (mean_val - cutoff)) & 
                            (sub[filter_metric] <= (mean_val + cutoff))
                        ]
                        cleaned_plot_df_list.append(filtered_sub)
                else:
                    cleaned_plot_df_list.append(sub)
            
            global_cleaned_df = pd.concat(cleaned_plot_df_list, ignore_index=True) if cleaned_plot_df_list else plot_df.copy()

            lines_removed = len(plot_df) - len(global_cleaned_df)
            if lines_removed > 0:
                st.toast(f"✂️ Pruned {lines_removed} outlier runs based on {filter_metric} ({sigma_multiplier}σ)!", icon="ℹ️")

            # --- LOOP THROUGH PLOT WINDOWS ---
            for i in range(st.session_state.num_plots):
                st.markdown(f"---")
                st.markdown(f"### Plot Window {i+1}")
                c1, c2 = st.columns(2)
                
                with c1:
                    selected_metric = st.selectbox(
                        f"Select Plotting Metric", 
                        available_metrics, 
                        key=f"metric_select_{i}",
                        index=available_metrics.index('PP-Gauss') if 'PP-Gauss' in available_metrics and i == 0 else i % len(available_metrics)
                    )
                with c2:
                    sort_order = st.selectbox(
                        "Sort Order",
                        ["Default (Name)", "Ascending (by Mean)", "Descending (by Mean)"],
                        key=f"sort_select_{i}"
                    )

                def process_and_render_data(data_to_plot, is_clean_view):
                    if view_mode == "Detailed (Show Every Run)":
                        temp_df = data_to_plot.copy()
                        
                        if sort_order != "Default (Name)":
                            group_means = temp_df.groupby('Sample Name')[selected_metric].mean()
                            ascending = True if "Ascending" in sort_order else False
                            sorted_names = group_means.sort_values(ascending=ascending).index
                            temp_df['Sample Name'] = pd.Categorical(temp_df['Sample Name'], categories=sorted_names, ordered=True)
                            temp_df = temp_df.sort_values(['Sample Name', 'Datestamp', 'Run'])
                        else:
                            temp_df = temp_df.sort_values(by=['Sample Name', 'Datestamp', 'Run'])

                        temp_df['Plot_X'] = (
                            "<span style='color:teal; font-weight:bold'>" + temp_df['Sample Name'].astype(str) + "</span><br>" + 
                            "<span style='color:gray'>" + temp_df['Datestamp'].astype(str) + "</span><br>" + 
                            "<span style='color:tomato'>Run: " + temp_df['Run'].astype(str) + "</span>"
                        )
                        
                        y_val = temp_df[selected_metric]
                        error_val = None
                        custom_data = temp_df[['Description']]
                        htemp = "<b>%{x}</b><br>Value: %{y}<br>Description: %{customdata[0]}<extra></extra>"
                        display_df = temp_df

                    else:
                        agg_results = data_to_plot.groupby('Sample Name')[available_metrics].agg(['mean', 'std']).reset_index()
                        agg_results.columns = [f"{col[0]}_{col[1]}" if col[1] else col[0] for col in agg_results.columns]
                        
                        if sort_order != "Default (Name)":
                            ascending = True if "Ascending" in sort_order else False
                            agg_results = agg_results.sort_values(by=f"{selected_metric}_mean", ascending=ascending)
                        
                        agg_results['Plot_X'] = "<span style='color:teal; font-weight:bold'>" + agg_results['Sample Name'] + "</span>"
                        
                        y_val = agg_results[f"{selected_metric}_mean"]
                        error_val = agg_results[f"{selected_metric}_std"]
                        custom_data = error_val
                        htemp = "<b>%{x}</b><br>Mean: %{y:.2f}<br>Std Dev: %{customdata:.2f}<extra></extra>"
                        display_df = agg_results

                    display_df['BarColor'] = display_df['Sample Name'].map(sticky_palette)

                    fig = go.Figure()
                    fig.add_trace(
                        go.Bar(
                            x=display_df['Plot_X'], 
                            y=y_val, 
                            marker_color=display_df['BarColor'],
                            text=y_val.round(2) if not y_val.isna().all() else "",
                            textposition='auto',
                            error_y=dict(type='data', array=error_val, visible=True) if view_mode == "Global (Mean & Std Dev)" else None,
                            customdata=custom_data,
                            hovertemplate=htemp,
                            textfont=dict(size=13, weight="bold")
                        )
                    )

                    title_lbl = f"{selected_metric} (Cleaned via {filter_metric} @ {sigma_multiplier}σ)" if is_clean_view else f"{selected_metric} (Unfiltered Original Data)"
                    
                    fig.update_layout(
                        title=dict(text=title_lbl, font=dict(size=18, weight="bold")),
                        height=480, 
                        yaxis_title=selected_metric, 
                        template="plotly_white", 
                        showlegend=False,
                        yaxis=dict(
                            range=[0, y_val.max() * 1.2] if not y_val.empty and not y_val.isna().all() else [0, 1],
                            title_font=dict(size=14, weight="bold"),
                            tickfont=dict(size=12)
                        ),
                        xaxis=dict(tickfont=dict(size=12)),
                        hoverlabel=dict(font_size=14)
                    )
                    fig.update_xaxes(tickangle=0)
                    return fig

                graph_col1, graph_col2 = st.columns(2)
                with graph_col1:
                    st.plotly_chart(process_and_render_data(plot_df, is_clean_view=False), width="stretch", key=f"chart_orig_{i}")
                with graph_col2:
                    st.plotly_chart(process_and_render_data(global_cleaned_df, is_clean_view=True), width="stretch", key=f"chart_clean_{i}")

            col_btn1, col_btn2 = st.columns([1, 4])
            with col_btn1:
                if st.button("➕ Add Plot", key="btn_add_plot"):
                    st.session_state.num_plots += 1
                    st.rerun()
            with col_btn2:
                if st.session_state.num_plots > 1:
                    if st.button("Reset Plots", key="btn_reset_plots"):
                        st.session_state.num_plots = 1
                        st.rerun()

            # --- SECTION 3: COHORT SUMMARY STATISTICS & CV% COMPARISON ---
            st.write("---")
            st.subheader("3. Cohort Summary Statistics Table & CV% Stability Delta")
            
            has_gauss = 'PP-Gauss' in plot_df.columns
            has_750 = 'PP-750' in plot_df.columns

            if (has_gauss or has_750) and not plot_df.empty:
                
                def generate_summary_metrics(dataframe, suffix):
                    agg_dict = {}
                    if has_gauss: agg_dict['PP-Gauss'] = ['count', 'mean', 'std']
                    if has_750: agg_dict['PP-750'] = ['count', 'mean', 'std']
                    
                    summary = dataframe.groupby('Sample Name').agg(agg_dict).reset_index()
                    summary.columns = [f"{col[0]}_{col[1]}" if col[1] else col[0] for col in summary.columns]
                    
                    final_cols = ['Sample Name']
                    if has_gauss:
                        summary[f'Total_Runs_Gauss_{suffix}'] = summary['PP-Gauss_count']
                        summary[f'Mean_Gauss_{suffix}'] = summary['PP-Gauss_mean']
                        summary[f'Std_Gauss_{suffix}'] = summary['PP-Gauss_std']
                        summary[f'CV%_Gauss_{suffix}'] = (summary[f'Std_Gauss_{suffix}'] / summary[f'Mean_Gauss_{suffix}']) * 100
                        final_cols.extend([f'Total_Runs_Gauss_{suffix}', f'Mean_Gauss_{suffix}', f'Std_Gauss_{suffix}', f'CV%_Gauss_{suffix}'])
                    
                    if has_750:
                        summary[f'Total_Runs_750_{suffix}'] = summary['PP-750_count']
                        summary[f'Mean_750_{suffix}'] = summary['PP-750_mean']
                        summary[f'Std_750_{suffix}'] = summary['PP-750_std']
                        summary[f'CV%_750_{suffix}'] = (summary[f'Std_750_{suffix}'] / summary[f'Mean_750_{suffix}']) * 100
                        final_cols.extend([f'Total_Runs_750_{suffix}', f'Mean_750_{suffix}', f'Std_750_{suffix}', f'CV%_750_{suffix}'])
                        
                    return summary[final_cols].fillna(0)

                orig_metrics = generate_summary_metrics(plot_df, "orig")
                clean_metrics = generate_summary_metrics(global_cleaned_df, "clean")
                merged_summary = pd.merge(orig_metrics, clean_metrics, on='Sample Name', how='outer').fillna(0)

                # --- CV% COMPARISON PLOT ---
                if has_gauss:
                    cv_fig = go.Figure()
                    cv_fig.add_trace(go.Bar(
                        x=merged_summary['Sample Name'],
                        y=merged_summary['CV%_Gauss_orig'],
                        name='Original Dataset PP-Gauss CV%',
                        marker_color='#ef553b',
                        text=merged_summary['CV%_Gauss_orig'].round(2).astype(str) + '%',
                        textposition='auto',
                        textfont=dict(size=12, weight="bold")
                    ))
                    cv_fig.add_trace(go.Bar(
                        x=merged_summary['Sample Name'],
                        y=merged_summary['CV%_Gauss_clean'],
                        name=f'Cleaned Dataset PP-Gauss CV% (via {filter_metric})',
                        marker_color='#636efa',
                        text=merged_summary['CV%_Gauss_clean'].round(2).astype(str) + '%',
                        textposition='auto',
                        textfont=dict(size=12, weight="bold")
                    ))
                    cv_fig.update_layout(
                        title=dict(text=f"Filtering Impact: PP-Gauss CV% Delta After Outlier Pruning", font=dict(size=18, weight="bold")),
                        xaxis_title="Sample Cohorts",
                        yaxis_title="Coefficient of Variation (CV %)",
                        barmode='group',
                        template='plotly_white',
                        height=420,
                        xaxis=dict(title_font=dict(size=14, weight="bold"), tickfont=dict(size=12)),
                        yaxis=dict(title_font=dict(size=14, weight="bold"), tickfont=dict(size=12)),
                        legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1, font=dict(size=12))
                    )
                    st.plotly_chart(cv_fig, width="stretch", key="cv_comparison_chart")

                # --- INTERACTIVE BREAKDOWN SUB-TABLE ---
                sum_c1, sum_c2 = st.columns([1, 3])
                with sum_c1:
                    summary_data_source = st.radio(
                        "Summary Metric Source for Table:",
                        ["Use Cleaned Dataset", "Use Original Dataset"],
                        key="summary_dataset_toggle"
                    )
                with sum_c2:
                    st.caption(f"Granular analytical cross-table showing tracking indices and stability benchmarks for **PP-Gauss** and **PP-750** simultaneously.")

                if summary_data_source == "Use Cleaned Dataset":
                    table_summary_view = clean_metrics.rename(columns={
                        'Total_Runs_Gauss_clean': 'Total_Runs',
                        'Mean_Gauss_clean': 'Mean_PP_Gauss', 'Std_Gauss_clean': 'Std_PP_Gauss', 'CV%_Gauss_clean': 'CV%_Gauss',
                        'Mean_750_clean': 'Mean_PP_750', 'Std_750_clean': 'Std_PP_750', 'CV%_750_clean': 'CV%_750'
                    })
                else:
                    table_summary_view = orig_metrics.rename(columns={
                        'Total_Runs_Gauss_orig': 'Total_Runs',
                        'Mean_Gauss_orig': 'Mean_PP_Gauss', 'Std_Gauss_orig': 'Std_PP_Gauss', 'CV%_Gauss_orig': 'CV%_Gauss',
                        'Mean_750_orig': 'Mean_PP_750', 'Std_750_orig': 'Std_PP_750', 'CV%_750_orig': 'CV%_750'
                    })

                st.dataframe(
                    table_summary_view[['Sample Name', 'Total_Runs', 'Mean_PP_Gauss', 'Std_PP_Gauss', 'CV%_Gauss', 'Mean_PP_750', 'Std_PP_750', 'CV%_750']],
                    column_config={
                        "Sample Name": st.column_config.TextColumn("Sample Name"),
                        "Total_Runs": st.column_config.NumberColumn("Total Runs", format="%d"),
                        "Mean_PP_Gauss": st.column_config.NumberColumn("Mean PP-Gauss", format="%.3f"),
                        "Std_PP_Gauss": st.column_config.NumberColumn("Std. Dev PP-Gauss", format="%.3f"),
                        "CV%_Gauss": st.column_config.NumberColumn("CV% (Gauss)", format="%.2f%%"),
                        "Mean_PP_750": st.column_config.NumberColumn("Mean PP-750", format="%.3f"),
                        "Std_PP_750": st.column_config.NumberColumn("Std. Dev PP-750", format="%.3f"),
                        "CV%_750": st.column_config.NumberColumn("CV% (750)", format="%.2f%%"),
                    },
                    hide_index=True,
                    width="stretch"
                )
            else:
                st.info("Metrics not found or dataset view empty. Summary table skipped.")

            # --- SECTION 4: PLOTTED DATA REFERENCE TABLE ---
            st.write("---")
            st.subheader("4. Plotted Data Reference Table")
            
            t_cfg1, t_cfg2 = st.columns([1, 3])
            with t_cfg1:
                table_view_mode = st.radio(
                    "Table Data Source:",
                    ["Show Cleaned Dataset", "Show Original Unfiltered Dataset"],
                    key="table_view_mode"
                )
            with t_cfg2:
                if table_view_mode == "Show Cleaned Dataset":
                    st.caption(f"Showing **{len(global_cleaned_df)} surviving rows**. Outliers based on variance in `{filter_metric}` have been filtered out.")
                else:
                    st.caption(f"Showing **{len(plot_df)} complete original rows**.")

            selected_table_source = global_cleaned_df if table_view_mode == "Show Cleaned Dataset" else plot_df
            table_df = selected_table_source.copy()
            
            if not table_df.empty:
                table_df['SAMPLE ID'] = table_df['Datestamp'].astype(str) + "_" + table_df['Sample Name'].astype(str) + "_" + table_df['Run'].astype(str)
                
                st.dataframe(
                    table_df[['Datestamp', 'Sample Name', 'Run', 'SAMPLE ID'] + available_metrics],
                    column_config={
                        "Datestamp": st.column_config.TextColumn("Date", width=100),
                        "Run": st.column_config.TextColumn("Run"), 
                        "SAMPLE ID": st.column_config.TextColumn("Combined ID", width=300),
                        **{m: st.column_config.NumberColumn(m, width=120, format="%.2f") for m in available_metrics}
                    },
                    hide_index=True,
                    width="stretch"
                )
            else:
                st.info("The selected dataset view is empty.")

        # --- SECTION 5: AUTOMATED DATA INSIGHTS ---
        st.divider()
        st.subheader("5. 🤖 Automated Data Insights")
        
        with st.expander("⚙️ AI Analysis & Custom Questions", expanded=True):
            i_col1, i_col2, i_col3 = st.columns([1, 1, 1])
            with i_col1:
                ai_metric = st.selectbox("Target Analysis Metric", available_metrics, key="ai_target_metric")
                lens = st.selectbox("Select Expert Lens", ["General Analyst", "Quality Control Specialist", "Root Cause Investigator"], key="ai_lens")
            with i_col2:
                custom_q = st.text_area("💬 Human Language Query", placeholder="e.g., 'Is there a correlation between NrDrops and the outliers in Sample B?'", key="ai_query")
            with i_col3:
                st.write(" ")
                st.write(" ")
                generate_btn = st.button("Generate Insights 🪄", width="stretch", key="btn_generate_insights")

        if generate_btn:
            with st.spinner("Analyzing..."):
                final_report = get_detailed_insights(plot_df, ai_metric, lens, custom_q)
                st.info(final_report)

        # --- SECTION 6: AUTOMATED POWERPOINT REPORT ENGINE ---
        st.divider()
        st.subheader("6. 📝 PowerPoint Summary Report Engine")
        st.caption("Generate a fully structured executive slide deck incorporating global metrics, cohort comparisons, and dedicated structural run slides per unique sample.")
        
        if st.button("Build PowerPoint Presentation Summary 🚀", width="stretch", key="btn_build_pptx"):
            with st.spinner("Rendering vector charts and drawing presentation canvas layers..."):
                try:
                    # 1. Initialize Blank PowerPoint File Layout (Widescreen 16:9 Aspect Ratio)
                    prs = Presentation()
                    prs.slide_width = Inches(13.333)
                    prs.slide_height = Inches(7.5)
                    blank_layout = prs.slide_layouts[6] 

                    # --- HELPER FUNCTION: ADD BASIC HEADER TO SLIDES ---
                    def add_slide_header(slide, title_text):
                        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.8))
                        tf = txBox.text_frame
                        tf.word_wrap = True
                        p = tf.paragraphs[0]
                        p.text = title_text
                        p.font.size = Pt(28)
                        p.font.bold = True
                        p.font.color.rgb = RGBColor(0, 51, 102)

                    # --- HELPER FUNCTION: BUILD MACRO OVERVIEW CHARTS ---
                    def build_macro_chart(dataframe, mode, is_clean, title):
                        temp = dataframe.copy()
                        # Enforce Ascending Sorting based on the group mean of PP-Gauss
                        group_means = temp.groupby('Sample Name')['PP-Gauss'].mean()
                        sorted_names = group_means.sort_values(ascending=True).index
                        temp['Sample Name'] = pd.Categorical(temp['Sample Name'], categories=sorted_names, ordered=True)
                        
                        fig = go.Figure()
                        if mode == "Detailed":
                            temp = temp.sort_values(['Sample Name', 'Datestamp', 'Run'])
                            temp['X_Label'] = temp['Sample Name'].astype(str) + "_" + temp['Run'].astype(str)
                            fig.add_trace(go.Bar(
                                x=temp['X_Label'], y=temp['PP-Gauss'],
                                marker_color=temp['Sample Name'].map(sticky_palette),
                                text=temp['PP-Gauss'].round(2), textposition='auto'
                            ))
                        else:
                            agg = temp.groupby('Sample Name')['PP-Gauss'].agg(['mean', 'std']).reset_index()
                            agg['Sample Name'] = pd.Categorical(agg['Sample Name'], categories=sorted_names, ordered=True)
                            agg = agg.sort_values('Sample Name')
                            fig.add_trace(go.Bar(
                                x=agg['Sample Name'], y=agg['mean'],
                                marker_color=agg['Sample Name'].map(sticky_palette),
                                error_y=dict(type='data', array=agg['std'], visible=True),
                                text=agg['mean'].round(2), textposition='auto'
                            ))
                            
                        fig.update_layout(
                            title=dict(text=title, font=dict(size=14, weight="bold")),
                            margin=dict(l=40, r=20, t=40, b=30), template="plotly_white", height=260,
                            yaxis=dict(title="PP-Gauss", title_font=dict(size=11), tickfont=dict(size=9)),
                            xaxis=dict(tickfont=dict(size=10, weight="bold")), showlegend=False
                        )
                        buf = BytesIO()
                        fig.write_image(buf, format="png", width=1150, height=240, scale=2)
                        buf.seek(0)
                        return buf

                    # --- SLIDE 1: DECK COVER TITLE SLIDE ---
                    slide1 = prs.slides.add_slide(blank_layout)
                    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(2.5))
                    tf = title_box.text_frame
                    p1 = tf.paragraphs[0]
                    p1.text = "Cohort Summary Report"
                    p1.font.size = Pt(44)
                    p1.font.bold = True
                    p1.font.color.rgb = RGBColor(0, 51, 102)
                    
                    p2 = tf.add_paragraph()
                    p2.text = f"Individual Runs Summary • Outliers Filtering - Baseline: {filter_metric} ({sigma_multiplier}σ)"
                    p2.font.size = Pt(18)
                    p2.font.color.rgb = RGBColor(102, 102, 102)

                    # --- SLIDE 2: NEW UNFILTERED GLOBAL PROFILE SLIDE ---
                    slide_macro_orig = prs.slides.add_slide(blank_layout)
                    add_slide_header(slide_macro_orig, "Global Analysis: Unfiltered Cohort Data")
                    
                    img_orig_det = build_macro_chart(plot_df, "Detailed", is_clean=False, title="1. All Individual Runs (Ascending by Group Mean)")
                    slide_macro_orig.shapes.add_picture(img_orig_det, Inches(0.666), Inches(1.3), width=Inches(12.0))
                    
                    img_orig_grp = build_macro_chart(plot_df, "Grouped", is_clean=False, title="2. Aggregated by Sample Name (Mean ± SD Dev)")
                    slide_macro_orig.shapes.add_picture(img_orig_grp, Inches(0.666), Inches(4.3), width=Inches(12.0))

                    # --- SLIDE 3: NEW CLEANED GLOBAL PROFILE SLIDE ---
                    slide_macro_clean = prs.slides.add_slide(blank_layout)
                    add_slide_header(slide_macro_clean, f"Global Analysis: Refined Clean Data ({sigma_multiplier}σ Outliers Pruned)")
                    
                    img_clean_det = build_macro_chart(global_cleaned_df, "Detailed", is_clean=True, title=f"3. Cleaned Individual Runs (via {filter_metric})")
                    slide_macro_clean.shapes.add_picture(img_clean_det, Inches(0.666), Inches(1.3), width=Inches(12.0))
                    
                    img_clean_grp = build_macro_chart(global_cleaned_df, "Grouped", is_clean=True, title=f"4. Cleaned Aggregated by Sample Name (Recalculated Mean ± SD Dev)")
                    slide_macro_clean.shapes.add_picture(img_clean_grp, Inches(0.666), Inches(4.3), width=Inches(12.0))

                    # --- SLIDE 4: COHORT STABILITY DELTA GRAPH ---
                    if has_gauss:
                        slide4 = prs.slides.add_slide(blank_layout)
                        add_slide_header(slide4, "Filtering Performance: PP-Gauss CV% Variations")
                        
                        img_stream = BytesIO()
                        cv_fig.write_image(img_stream, format="png", width=1100, height=500, scale=2)
                        img_stream.seek(0)
                        slide4.shapes.add_picture(img_stream, Inches(0.666), Inches(1.3), width=Inches(12.0))

                    # --- SLIDE 5: COHORT DATA TABLE SUMMARY ---
                    slide5 = prs.slides.add_slide(blank_layout)
                    add_slide_header(slide5, f"Cohort Parameter Analytics View ({summary_data_source})")
                    
                    ppt_table_df = table_summary_view[['Sample Name', 'Total_Runs', 'Mean_PP_Gauss', 'Std_PP_Gauss', 'CV%_Gauss', 'Mean_PP_750', 'Std_PP_750', 'CV%_750']].reset_index(drop=True)
                    
                    rows, cols = len(ppt_table_df) + 1, len(ppt_table_df.columns)
                    left, top, width, height = Inches(0.5), Inches(1.5), Inches(12.333), Inches(0.4 * rows)
                    table_shape = slide5.shapes.add_table(rows, cols, left, top, width, height)
                    table = table_shape.table

                    headers_display = ["Sample Name", "Total Runs", "Mean Gauss", "Std Dev Gauss", "CV% Gauss", "Mean 750", "Std Dev 750", "CV% 750"]
                    for col_idx, text in enumerate(headers_display):
                        cell = table.cell(0, col_idx)
                        cell.text = text
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
                        p = cell.text_frame.paragraphs[0]
                        p.font.bold = True
                        p.font.color.rgb = RGBColor(255, 255, 255)
                        p.font.size = Pt(12)
                        p.alignment = PP_ALIGN.CENTER

                    for sequential_row_idx in range(len(ppt_table_df)):
                        row_series = ppt_table_df.iloc[sequential_row_idx]
                        for col_idx in range(len(row_series)):
                            cell = table.cell(sequential_row_idx + 1, col_idx)
                            value = row_series.iloc[col_idx]
                            
                            col_name = ppt_table_df.columns[col_idx]
                            if isinstance(value, float):
                                cell.text = f"{value:.3f}%" if "CV%" in col_name else f"{value:.3f}"
                            else:
                                cell.text = str(value)
                                
                            p = cell.text_frame.paragraphs[0]
                            p.font.size = Pt(11)
                            p.alignment = PP_ALIGN.CENTER

                    # --- SLIDES 6+: INDIVIDUAL COHORT RUN TIMELINES ---
                    def build_micro_timeline(dataframe, sample_name, y_metric, title_text):
                        sub_sub = dataframe[dataframe['Sample Name'] == sample_name].sort_values(by=['Datestamp', 'Run'])
                        sub_sub['X_Label'] = "Run: " + sub_sub['Run'].astype(str)
                        
                        fig_micro = go.Figure(data=[
                            go.Bar(
                                x=sub_sub['X_Label'],
                                y=sub_sub[y_metric],
                                marker_color=sticky_palette.get(sample_name, '#636efa'),
                                text=sub_sub[y_metric].round(2) if not sub_sub[y_metric].isna().all() else "",
                                textposition='auto'
                            )
                        ])
                        fig_micro.update_layout(
                            title=dict(text=title_text, font=dict(size=14, weight="bold")),
                            margin=dict(l=40, r=20, t=40, b=30),
                            template="plotly_white",
                            height=250,
                            yaxis=dict(title=y_metric, title_font=dict(size=10), tickfont=dict(size=9)),
                            xaxis=dict(tickfont=dict(size=10))
                        )
                        img_buf = BytesIO()
                        fig_micro.write_image(img_buf, format="png", width=1100, height=220, scale=2)
                        img_buf.seek(0)
                        return img_buf

                    unique_cohorts = sorted(plot_df['Sample Name'].unique())
                    
                    for cohort in unique_cohorts:
                        slide_c = prs.slides.add_slide(blank_layout)
                        add_slide_header(slide_c, f"Sample Run Analysis: {cohort} (Outliers Removed)")
                        
                        if 'PP-Gauss' in plot_df.columns:
                            img_g = build_micro_timeline(global_cleaned_df, cohort, 'PP-Gauss', "PP-Gauss")
                            slide_c.shapes.add_picture(img_g, Inches(0.666), Inches(1.2), width=Inches(12.0))
                            
                        if 'PP-750' in plot_df.columns:
                            img_750 = build_micro_timeline(global_cleaned_df, cohort, 'PP-750', "PP-750")
                            slide_c.shapes.add_picture(img_750, Inches(0.666), Inches(3.2), width=Inches(12.0))
                            
                        if 'Barcode av.' in plot_df.columns:
                            img_bav = build_micro_timeline(global_cleaned_df, cohort, 'Barcode av.', "Barcode average")
                            slide_c.shapes.add_picture(img_bav, Inches(0.666), Inches(5.2), width=Inches(12.0))

                    # Save and Deliver Presentation File Stream
                    ppt_out = BytesIO()
                    prs.save(ppt_out)
                    ppt_out.seek(0)
                    
                    st.success("PowerPoint compilation complete! Click below to download.")
                    st.download_button(
                        label="📥 Download PowerPoint Report (.pptx)",
                        data=ppt_out,
                        file_name="MJFF_Instrumentation_Summary_Report.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                except Exception as e:
                    st.error(f"Failed to generate presentation deck: {str(e)}")